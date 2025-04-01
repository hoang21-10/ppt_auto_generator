import streamlit as st
import os
import pptx
from pptx.util import Pt
from pptx.dml.color import RGBColor
from docx import Document
from dotenv import load_dotenv
import google.generativeai as genai
import requests
from bs4 import BeautifulSoup
from firebase_config import auth  # Import Firebase Auth

# Load API Key từ .env
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=GOOGLE_API_KEY)
MODEL_NAME = "gemini-1.5-flash"
model = genai.GenerativeModel(MODEL_NAME)

TEMP_DIR = "temp"
os.makedirs(TEMP_DIR, exist_ok=True)

# ✅ Hàm đăng nhập
def login():
    st.subheader("🔑 Đăng nhập")
    email = st.text_input("📧 Email")
    password = st.text_input("🔒 Mật khẩu", type="password")

    if st.button("Đăng nhập"):
        try:
            user = auth.sign_in_with_email_and_password(email, password)
            st.session_state["user"] = user
            st.success("✅ Đăng nhập thành công!")
            st.experimental_rerun()
        except:
            st.error("❌ Sai email hoặc mật khẩu!")

# ✅ Hàm đăng ký tài khoản
def register():
    st.subheader("📝 Đăng ký")
    email = st.text_input("📧 Email", key="reg_email")
    password = st.text_input("🔒 Mật khẩu", type="password", key="reg_password")

    if st.button("Đăng ký"):
        try:
            auth.create_user_with_email_and_password(email, password)
            st.success("✅ Đăng ký thành công! Hãy đăng nhập.")
        except:
            st.error("❌ Email đã tồn tại hoặc lỗi khác.")

# ✅ Hàm đăng xuất
def logout():
    if st.button("Đăng xuất"):
        del st.session_state["user"]
        st.rerun()

# ✅ Hàm đọc file TXT
def read_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as file:
        return file.read()

# ✅ Hàm đọc file DOCX (bao gồm cả bảng)
def read_docx(file_path):
    doc = Document(file_path)
    content = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            content.append(text)
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            table_data.append([cell.text.strip() for cell in row.cells])
        content.append("\n".join([" | ".join(row) for row in table_data]))
    return "\n".join(content)

# ✅ Hàm lấy nội dung từ URL
def fetch_article_content(url):
    try:
        response = requests.get(url, headers={"User-Agent": "Mozilla/5.0"})
        if response.status_code == 200:
            soup = BeautifulSoup(response.text, "html.parser")
            paragraphs = soup.find_all("p")
            content = "\n".join([p.get_text() for p in paragraphs])
            return content
        return "Không thể lấy nội dung từ URL."
    except Exception as e:
        return f"Lỗi: {e}"

# ✅ Hàm chia nội dung thành các slide
def split_content_into_slides(content, max_chars=300):
    slides = []
    paragraphs = content.split("\n")
    for para in paragraphs:
        para = para.strip()
        if para:
            sentences = para.split(". ")
            slide_text = ""
            for sentence in sentences:
                if len(slide_text) + len(sentence) > max_chars:
                    slides.append(("Nội dung", slide_text))
                    slide_text = sentence + "."
                else:
                    slide_text += " " + sentence if slide_text else sentence
            if slide_text:
                slides.append(("Nội dung", slide_text))
    return slides

# ✅ Hàm tạo file PowerPoint
def create_presentation(title, slides, font_size, font_color_rgb):
    ppt = pptx.Presentation()
    title_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    title_slide.shapes.title.text = title
    
    for slide_title, slide_text in slides:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        slide.shapes.title.text = slide_title
        text_box = slide.shapes.placeholders[1]
        text_frame = text_box.text_frame
        text_frame.clear()
        
        paragraphs = slide_text.split(". ")
        for para in paragraphs:
            p = text_frame.add_paragraph()
            p.text = para + "."
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(*font_color_rgb)
    
    ppt_path = os.path.join(TEMP_DIR, "presentation.pptx")
    ppt.save(ppt_path)
    return ppt_path

# ✅ Giao diện chính của ứng dụng
def main():
    st.title("📑 Chuyển đổi TXT/DOCX/URL sang PPT")
    
    if "user" not in st.session_state:
        login()
        st.markdown("---")
        register()
        return
    
    logout()
    
    uploaded_file = st.file_uploader("📌 Tải lên file", type=["txt", "docx"])
    url_input = st.text_input("🌐 Nhập URL bài viết")
    
    font_size = st.sidebar.slider("🖋 Kích thước chữ", 12, 48, 24)
    font_color = st.sidebar.color_picker("🎨 Màu chữ", "#000000")
    font_color_rgb = tuple(int(font_color[i:i+2], 16) for i in (1, 3, 5))
    
    content = ""
    if uploaded_file:
        file_path = os.path.join(TEMP_DIR, uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        content = read_txt(file_path) if uploaded_file.name.endswith(".txt") else read_docx(file_path)
    elif url_input:
        content = fetch_article_content(url_input)
    
    if content:
        slides = split_content_into_slides(content)
        st.subheader("📌 Xem trước nội dung")
        for idx, (title, text) in enumerate(slides):
            with st.expander(f"🔹 Slide {idx + 1}: {title}"):
                st.write(text)
        
        if st.button("✨ Xuất PowerPoint"):
            ppt_path = create_presentation("Bài thuyết trình", slides, font_size, font_color_rgb)
            st.success("✅ Đã tạo xong PowerPoint!")
            with open(ppt_path, "rb") as file:
                st.download_button("📥 Tải về PowerPoint", file, "presentation.pptx")

if __name__ == "__main__":
    main()